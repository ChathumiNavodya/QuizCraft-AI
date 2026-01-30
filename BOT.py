import streamlit as st
import sqlite3
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.dates as mdates

import PyPDF2
import os
from datetime import datetime
import hashlib
import json
import time
import re
import secrets
from typing import Tuple

from docx import Document
from striprtf.striprtf import rtf_to_text
from odf.opendocument import load
from odf import teletype
from pptx import Presentation

import google.generativeai as genai


# ============================================================
# App Branding
# ============================================================
APP_NAME = "QuizCraft AI"
APP_ICON = "🧠"
st.set_page_config(
    page_title=APP_NAME,
    page_icon=APP_ICON,
    layout="wide",
    initial_sidebar_state="expanded"
)
SCHEMA_VERSION = "2026-01-30-v2-fixed"


# ============================================================
# Paths
# ============================================================
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DB_PATH = os.path.join(BASE_DIR, "quiz_performance.db")


# ============================================================
# Dark Theme (Polished) 🌙
# ============================================================
st.markdown(
    """
    <style>
      @import url('https://fonts.googleapis.com/css2?family=Outfit:wght@300;400;600;700;800&display=swap');

      :root {
        --bg0: #070b14;
        --bg1: #0b1220;
        --bg2: #0f1b2d;
        --text0: #e8f0ff;
        --text1: #c8d4f0;
        --muted: #8aa0c8;
        --primary-teal: #22d3ee;
        --deep-teal: #0891b2;
        --accent-orange: #fb923c;
        --accent-coral: #fb7185;
        --accent-yellow: #fbbf24;
        --border: rgba(255,255,255,0.12);
        --shadow: 0 18px 60px rgba(0,0,0,0.45);
      }

      * { font-family: 'Outfit', -apple-system, BlinkMacSystemFont, sans-serif; }

      .stApp {
        background: radial-gradient(1200px 600px at 20% 0%, rgba(34,211,238,0.14), transparent 60%),
                    radial-gradient(900px 500px at 95% 15%, rgba(251,146,60,0.10), transparent 55%),
                    linear-gradient(180deg, var(--bg0) 0%, var(--bg1) 45%, var(--bg2) 100%);
        background-attachment: fixed;
        color: var(--text0);
      }

      .block-container {
        padding-top: 1.2rem;
        padding-bottom: 2.5rem;
        max-width: 1400px;
      }

      [data-testid="stSidebar"] {
        min-width: 320px;
        max-width: 380px;
        background: linear-gradient(180deg, #070b14 0%, #0b1220 60%, #0f1b2d 100%);
        box-shadow: 6px 0 28px rgba(0,0,0,0.55);
        border-right: 1px solid var(--border);
      }
      [data-testid="stSidebar"] * { color: var(--text0) !important; }

      .hero-title {
        font-size: 3.4rem;
        font-weight: 800;
        background: linear-gradient(135deg, var(--accent-yellow) 0%, var(--accent-orange) 35%, var(--accent-coral) 70%, var(--primary-teal) 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin-bottom: 0.4rem;
        line-height: 1.15;
        text-align: center;
        letter-spacing: -1px;
      }
      .hero-subtitle {
        font-size: 1.15rem;
        color: var(--text1);
        font-weight: 400;
        margin-bottom: 1.6rem;
        text-align: center;
      }

      .qc-card {
        padding: 2.0rem;
        border-radius: 22px;
        border: 1px solid var(--border);
        background: linear-gradient(135deg, rgba(255,255,255,0.07), rgba(255,255,255,0.04));
        backdrop-filter: blur(14px);
        box-shadow: var(--shadow);
        margin-bottom: 1.2rem;
      }

      .feature-card {
        background: linear-gradient(135deg, rgba(255,255,255,0.08), rgba(255,255,255,0.04));
        padding: 1.8rem 1.4rem;
        border-radius: 20px;
        text-align: center;
        height: 100%;
        border: 1px solid var(--border);
        box-shadow: 0 10px 30px rgba(0,0,0,0.35);
        position: relative;
        overflow: hidden;
      }
      .feature-card::before {
        content: '';
        position: absolute;
        top: 0; left: 0; width: 100%; height: 4px;
        background: linear-gradient(90deg, var(--accent-yellow), var(--accent-orange), var(--primary-teal));
        transform: scaleX(0);
        transition: transform 0.35s ease;
      }
      .feature-card:hover::before { transform: scaleX(1); }
      .feature-icon { font-size: 3.2rem; margin-bottom: 1.0rem; display:inline-block; }
      .feature-title { font-size: 1.15rem; font-weight: 800; color: var(--primary-teal); margin-bottom: 0.4rem; }
      .feature-desc { color: var(--text1); font-size: 0.98rem; line-height: 1.5; }

      .stTextInput>div>div>input,
      .stNumberInput>div>div>input,
      .stSelectbox>div>div>select,
      textarea {
        border-radius: 14px !important;
        border: 1px solid rgba(255,255,255,0.14) !important;
        background: rgba(255,255,255,0.06) !important;
        color: var(--text0) !important;
      }

      .stButton>button {
        border-radius: 14px;
        font-weight: 700;
        padding: 0.85rem 1.05rem;
        border: 1px solid rgba(34,211,238,0.28);
        background: linear-gradient(135deg, rgba(34,211,238,0.26), rgba(8,145,178,0.22));
        color: var(--text0);
        box-shadow: 0 10px 26px rgba(0,0,0,0.35);
        width: 100%;
        transition: transform 0.10s ease, filter 0.20s ease;
      }
      .stButton>button:hover { filter: brightness(1.15); transform: translateY(-2px); }
      .stButton>button:active { transform: translateY(0px); }

      [data-testid="stFileUploader"] {
        background: linear-gradient(135deg, rgba(34,211,238,0.08), rgba(251,146,60,0.04));
        border-radius: 20px;
        padding: 2.2rem 1.8rem;
        border: 2px dashed rgba(34,211,238,0.35);
      }

      .quiz-question {
        background: linear-gradient(135deg, rgba(255,255,255,0.075), rgba(255,255,255,0.04));
        padding: 1.8rem;
        border-radius: 18px;
        margin-bottom: 1.1rem;
        border: 1px solid var(--border);
        border-left: 6px solid rgba(34,211,238,0.7);
        box-shadow: 0 14px 34px rgba(0,0,0,0.35);
      }

      hr {
        margin: 2.0rem 0;
        border: none;
        height: 2px;
        background: linear-gradient(90deg, transparent, rgba(34,211,238,0.35), rgba(251,146,60,0.25), transparent);
        border-radius: 3px;
      }

      @media (max-width: 768px) {
        .hero-title { font-size: 2.2rem !important; }
        .qc-card { padding: 1.4rem; }
      }
    </style>
    """,
    unsafe_allow_html=True
)


# ============================================================
# Helpers
# ============================================================
def compute_hash(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()

def hash_password(password: str, salt: str) -> str:
    dk = hashlib.pbkdf2_hmac("sha256", password.encode("utf-8"), salt.encode("utf-8"), 200_000)
    return dk.hex()

def new_salt() -> str:
    return secrets.token_hex(16)


# ============================================================
# Gemini Initialization (Fixed - Auto-detect working model)
# ============================================================
def initialize_llm():
    """Initialize Gemini model with auto-detection of available models"""
    api_key = st.secrets.get("GOOGLE_API_KEY", None) or os.getenv("GOOGLE_API_KEY")
    
    if not api_key:
        st.error("⚠️ Missing GOOGLE_API_KEY")
        st.info("Add GOOGLE_API_KEY to Streamlit Cloud → Settings → Secrets")
        return None

    try:
        genai.configure(api_key=api_key)
        
        # Try preferred models first
        preferred_models = [
            "gemini-2.0-flash-exp",
            "gemini-1.5-flash",
            "gemini-1.5-flash-001",
            "gemini-1.5-pro",
            "gemini-1.5-pro-001",
        ]
        
        for model_name in preferred_models:
            try:
                model = genai.GenerativeModel(model_name)
                # Test the model with a simple prompt
                test_response = model.generate_content("Hello")
                if test_response:
                    return model
            except Exception:
                continue
        
        # If none of the preferred models work, auto-discover
        available_models = []
        for m in genai.list_models():
            if "generateContent" in (getattr(m, "supported_generation_methods", None) or []):
                available_models.append(m.name)
        
        if available_models:
            model_name = available_models[0].replace("models/", "")
            return genai.GenerativeModel(model_name)
        
        st.error("No compatible Gemini models found for your API key")
        return None
        
    except Exception as e:
        st.error(f"Failed to initialize Gemini: {e}")
        return None


# ============================================================
# Database Setup
# ============================================================
@st.cache_resource
def init_db(schema_version: str):
    conn = sqlite3.connect(DB_PATH, check_same_thread=False)
    conn.execute("PRAGMA foreign_keys = ON;")

    conn.executescript("""
        CREATE TABLE IF NOT EXISTS users (
            user_id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE,
            email TEXT UNIQUE,
            salt TEXT,
            password_hash TEXT,
            created_at TEXT,
            email_verified INTEGER DEFAULT 1
        );

        CREATE TABLE IF NOT EXISTS modules (
            module_id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            module_name TEXT,
            created_at TEXT,
            UNIQUE(user_id, module_name),
            FOREIGN KEY (user_id) REFERENCES users(user_id)
        );

        CREATE TABLE IF NOT EXISTS quizzes (
            quiz_id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            module_id INTEGER,
            timestamp TEXT,
            notes_hash TEXT,
            difficulty TEXT,
            FOREIGN KEY (user_id) REFERENCES users(user_id),
            FOREIGN KEY (module_id) REFERENCES modules(module_id)
        );

        CREATE TABLE IF NOT EXISTS questions (
            question_id INTEGER PRIMARY KEY AUTOINCREMENT,
            quiz_id INTEGER,
            question_text TEXT,
            options TEXT,
            correct_answer INTEGER,
            explanation TEXT,
            FOREIGN KEY (quiz_id) REFERENCES quizzes(quiz_id)
        );

        CREATE TABLE IF NOT EXISTS user_answers (
            answer_id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            quiz_id INTEGER,
            question_id INTEGER,
            user_answer INTEGER,
            is_correct INTEGER,
            FOREIGN KEY (user_id) REFERENCES users(user_id),
            FOREIGN KEY (quiz_id) REFERENCES quizzes(quiz_id),
            FOREIGN KEY (question_id) REFERENCES questions(question_id)
        );

        CREATE TABLE IF NOT EXISTS performance (
            performance_id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            module_id INTEGER,
            timestamp TEXT,
            score INTEGER,
            total INTEGER,
            percentage REAL,
            FOREIGN KEY (user_id) REFERENCES users(user_id),
            FOREIGN KEY (module_id) REFERENCES modules(module_id)
        );
    """)
    conn.commit()
    return conn


# ============================================================
# Auth Functions
# ============================================================
def create_user(conn, username: str, email: str, password: str) -> Tuple[bool, str]:
    username = (username or "").strip()
    email = (email or "").strip().lower()
    password = (password or "").strip()

    if not username or not password:
        return False, "Username and password required."

    email_regex = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
    if not re.match(email_regex, email):
        return False, "Invalid email format."

    if len(password) < 6 or not re.search(r"\d", password):
        return False, "Password must be at least 6 characters and include a number."

    try:
        salt = new_salt()
        ph = hash_password(password, salt)

        conn.execute(
            "INSERT INTO users (username, email, salt, password_hash, created_at, email_verified) VALUES (?, ?, ?, ?, ?, ?)",
            (username, email, salt, ph, datetime.now().isoformat(timespec="seconds"), 1)
        )
        conn.commit()
        return True, "Account created successfully! You can now login."

    except sqlite3.IntegrityError as e:
        if "username" in str(e).lower():
            return False, "Username already exists."
        if "email" in str(e).lower():
            return False, "Email already registered."
        return False, f"Registration error: {e}"


def login_user(conn, username_or_email: str, password: str) -> Tuple[bool, str, int, str]:
    c = conn.cursor()
    identifier = (username_or_email or "").strip()
    password = (password or "").strip()

    if not identifier or not password:
        return False, "Enter username/email and password.", -1, ""

    c.execute(
        """
        SELECT user_id, salt, password_hash, email_verified, username
        FROM users
        WHERE LOWER(username) = LOWER(?)
           OR LOWER(email) = LOWER(?)
        """,
        (identifier, identifier)
    )
    row = c.fetchone()

    if not row:
        return False, "Account not found. Please register first.", -1, ""

    user_id, salt, ph, email_verified, username = row

    if hash_password(password, salt) != ph:
        return False, "Incorrect password.", -1, ""

    return True, f"Welcome back, {username}!", user_id, username


# ============================================================
# Module Functions
# ============================================================
def get_or_create_module(conn, user_id: int, module_name: str):
    module_name = (module_name or "").strip()
    if not module_name:
        return None

    c = conn.cursor()
    c.execute("SELECT module_id FROM modules WHERE user_id = ? AND module_name = ?", (user_id, module_name))
    row = c.fetchone()
    if row:
        return row[0]

    c.execute(
        "INSERT INTO modules (user_id, module_name, created_at) VALUES (?, ?, ?)",
        (user_id, module_name, datetime.now().isoformat(timespec="seconds"))
    )
    conn.commit()
    return c.lastrowid

def list_modules(conn, user_id: int):
    c = conn.cursor()
    c.execute("SELECT module_id, module_name FROM modules WHERE user_id = ? ORDER BY module_name ASC", (user_id,))
    return c.fetchall()


# ============================================================
# Quiz Functions
# ============================================================
def save_quiz(conn, user_id: int, module_id: int, notes_hash: str, quiz_data: dict, difficulty: str):
    c = conn.cursor()
    timestamp = datetime.now().isoformat(timespec="seconds")

    c.execute(
        "INSERT INTO quizzes (user_id, module_id, timestamp, notes_hash, difficulty) VALUES (?, ?, ?, ?, ?)",
        (user_id, module_id, timestamp, notes_hash, difficulty),
    )
    quiz_id = c.lastrowid

    for q in quiz_data["questions"]:
        c.execute(
            "INSERT INTO questions (quiz_id, question_text, options, correct_answer, explanation) VALUES (?, ?, ?, ?, ?)",
            (quiz_id, q["question"], json.dumps(q["options"]), int(q["answer"]), q.get("explanation", "")),
        )

    conn.commit()
    return quiz_id

def list_quiz_ids_for_module(conn, user_id: int, module_id: int, limit: int = 50):
    return conn.execute(
        "SELECT quiz_id, timestamp, difficulty FROM quizzes WHERE user_id=? AND module_id=? ORDER BY quiz_id DESC LIMIT ?",
        (user_id, module_id, limit)
    ).fetchall()

def load_quiz(conn, quiz_id: int):
    qrow = conn.execute("SELECT quiz_id, difficulty, timestamp FROM quizzes WHERE quiz_id=?", (quiz_id,)).fetchone()
    if not qrow:
        return None

    qdf = pd.read_sql_query(
        "SELECT question_id, question_text, options, correct_answer, explanation FROM questions WHERE quiz_id=? ORDER BY question_id ASC",
        conn,
        params=(quiz_id,)
    )

    questions = []
    for _, r in qdf.iterrows():
        questions.append({
            "question_id": int(r["question_id"]),
            "question": r["question_text"],
            "options": json.loads(r["options"]),
            "answer": int(r["correct_answer"]),
            "explanation": r["explanation"] or ""
        })

    return {"quiz_id": qrow[0], "difficulty": qrow[1], "timestamp": qrow[2], "questions": questions}

def save_user_answers_and_performance(conn, user_id: int, module_id: int, quiz: dict, user_choices: dict):
    quiz_id = quiz["quiz_id"]
    total = len(quiz["questions"])
    score = 0

    conn.execute("DELETE FROM user_answers WHERE user_id=? AND quiz_id=?", (user_id, quiz_id))

    for q in quiz["questions"]:
        qid = q["question_id"]
        correct = int(q["answer"])
        chosen = int(user_choices.get(qid, -1))
        is_correct = 1 if chosen == correct else 0
        score += is_correct

        conn.execute(
            "INSERT INTO user_answers (user_id, quiz_id, question_id, user_answer, is_correct) VALUES (?, ?, ?, ?, ?)",
            (user_id, quiz_id, qid, chosen, is_correct)
        )

    percentage = (score / total) * 100 if total else 0
    conn.execute(
        "INSERT INTO performance (user_id, module_id, timestamp, score, total, percentage) VALUES (?, ?, ?, ?, ?, ?)",
        (user_id, module_id, datetime.now().isoformat(timespec="seconds"), score, total, percentage)
    )
    conn.commit()
    return score, total, percentage

def fetch_performance(conn, user_id: int, module_id: int):
    df = pd.read_sql_query(
        "SELECT timestamp, score, total, percentage FROM performance WHERE user_id=? AND module_id=? ORDER BY timestamp ASC",
        conn,
        params=(user_id, module_id),
    )
    if not df.empty:
        df["timestamp"] = pd.to_datetime(df["timestamp"])
    return df


# ============================================================
# File Extraction
# ============================================================
def extract_notes(file) -> Tuple[str, str]:
    try:
        name = file.name.lower()

        if name.endswith(".pdf"):
            reader = PyPDF2.PdfReader(file)
            text = "\n".join([(p.extract_text() or "") for p in reader.pages]).strip()
            if len(text) < 50:
                return "", "PDF appears to be scanned or empty. Please use a text-based PDF."
            return text, ""

        if name.endswith(".docx"):
            doc = Document(file)
            text = "\n".join([p.text for p in doc.paragraphs]).strip()
            return text, ""

        if name.endswith(".rtf"):
            raw = file.read()
            if isinstance(raw, bytes):
                raw = raw.decode("utf-8", errors="ignore")
            text = rtf_to_text(raw).strip()
            return text, ""

        if name.endswith(".odt"):
            doc = load(file)
            text = teletype.extractText(doc.text).strip()
            return text, ""

        if name.endswith(".pptx"):
            prs = Presentation(file)
            all_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text.append(shape.text)
            return "\n".join(all_text).strip(), ""

        # Default text file
        raw = file.read()
        if isinstance(raw, bytes):
            raw = raw.decode("utf-8", errors="ignore")
        return raw.strip(), ""

    except Exception as e:
        return "", f"Error extracting text: {str(e)}"


# ============================================================
# Quiz Generation (Fixed)
# ============================================================
def validate_quiz_json(data):
    if not isinstance(data, dict) or "questions" not in data:
        return False
    if not isinstance(data["questions"], list):
        return False
    
    for q in data["questions"]:
        if not isinstance(q, dict):
            return False
        if not {"question", "options", "answer"}.issubset(q.keys()):
            return False
        if not isinstance(q["options"], list) or len(q["options"]) != 4:
            return False
        try:
            ans = int(q["answer"])
            if not (0 <= ans <= 3):
                return False
        except:
            return False
    
    return True

def _extract_json(text: str) -> str:
    """Extract JSON from markdown code blocks or plain text"""
    text = (text or "").strip()
    # Remove markdown code blocks
    text = re.sub(r"^```(?:json)?\s*", "", text)
    text = re.sub(r"\s*```$", "", text)
    # Find JSON object
    match = re.search(r"\{.*\}", text, flags=re.DOTALL)
    return match.group(0) if match else text

def generate_quiz_once(llm, notes_text: str, num_questions: int, difficulty: str):
    """Generate quiz with single attempt"""
    prompt = f"""You are an expert quiz generator. Create {num_questions} multiple-choice questions based ONLY on the provided notes.

Difficulty level: {difficulty}

Return ONLY valid JSON in this exact format (no markdown, no extra text):
{{"questions": [{{"question":"Question text here","options":["Option A","Option B","Option C","Option D"],"answer":0,"explanation":"Brief explanation"}}]}}

Rules:
- "answer" must be 0, 1, 2, or 3 (index of correct option)
- Include exactly 4 options per question
- Keep explanations brief and clear
- Base questions only on the provided notes

NOTES:
{notes_text[:4000]}
"""

    response = llm.generate_content(prompt)
    raw_text = getattr(response, "text", "") or ""
    
    # Extract JSON
    json_str = _extract_json(raw_text)
    
    try:
        data = json.loads(json_str)
    except json.JSONDecodeError as e:
        raise ValueError(f"Invalid JSON from model: {e}")
    
    if not validate_quiz_json(data):
        raise ValueError("Model returned invalid quiz structure")
    
    # Ensure all fields are correct type
    for q in data["questions"]:
        q["answer"] = int(q["answer"])
        if "explanation" not in q:
            q["explanation"] = ""
    
    return data

def generate_quiz_with_retry(llm, notes_text: str, num_questions: int, difficulty: str, max_attempts: int = 3):
    """Generate quiz with retry logic"""
    last_error = None
    
    for attempt in range(1, max_attempts + 1):
        try:
            return generate_quiz_once(llm, notes_text, num_questions, difficulty)
        except Exception as e:
            last_error = e
            if attempt < max_attempts:
                time.sleep(1.5)
            else:
                raise ValueError(f"Failed after {max_attempts} attempts: {last_error}")


# ============================================================
# Plotting
# ============================================================
def plot_progress(df: pd.DataFrame):
    fig, ax = plt.subplots(figsize=(12, 5))
    ax.plot(df["timestamp"], df["percentage"], marker="o", linestyle="-", linewidth=2.5, markersize=8, color='#22d3ee')
    ax.fill_between(df["timestamp"], df["percentage"], alpha=0.15, color='#22d3ee')
    ax.set_title("Performance Over Time", fontsize=16, fontweight='bold', pad=20, color='#e8f0ff')
    ax.set_xlabel("Date", fontsize=12, color='#c8d4f0')
    ax.set_ylabel("Score (%)", fontsize=12, color='#c8d4f0')
    ax.set_ylim(0, 105)
    ax.grid(True, alpha=0.2, linestyle='--', color='#8aa0c8')
    
    # Style the plot
    ax.set_facecolor('#0b1220')
    fig.patch.set_facecolor('#0b1220')
    ax.tick_params(colors='#c8d4f0')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#8aa0c8')
    ax.spines['left'].set_color('#8aa0c8')

    # Format dates
    locator = mdates.AutoDateLocator(minticks=3, maxticks=7)
    formatter = mdates.ConciseDateFormatter(locator)
    ax.xaxis.set_major_locator(locator)
    ax.xaxis.set_major_formatter(formatter)
    
    fig.tight_layout()
    return fig


# ============================================================
# UI Components
# ============================================================
def render_hero():
    st.markdown(f'<div class="hero-title">{APP_ICON} {APP_NAME}</div>', unsafe_allow_html=True)
    st.markdown('<div class="hero-subtitle">Transform your study notes into interactive quizzes powered by AI</div>', unsafe_allow_html=True)

def render_features():
    col1, col2, col3 = st.columns(3)
    with col1:
        st.markdown("""
        <div class="feature-card">
            <span class="feature-icon">📝</span>
            <div class="feature-title">Smart Quiz Generation</div>
            <div class="feature-desc">Upload notes and generate quizzes automatically using Google Gemini AI</div>
        </div>
        """, unsafe_allow_html=True)
    with col2:
        st.markdown("""
        <div class="feature-card">
            <span class="feature-icon">✅</span>
            <div class="feature-title">Interactive Testing</div>
            <div class="feature-desc">Take quizzes with instant scoring and detailed explanations</div>
        </div>
        """, unsafe_allow_html=True)
    with col3:
        st.markdown("""
        <div class="feature-card">
            <span class="feature-icon">📊</span>
            <div class="feature-title">Progress Tracking</div>
            <div class="feature-desc">Track your performance over time with visual analytics</div>
        </div>
        """, unsafe_allow_html=True)


# ============================================================
# Main Application
# ============================================================
def main():
    # Initialize session state
    if "active_tab" not in st.session_state:
        st.session_state["active_tab"] = "notes"

    # Initialize database and LLM
    conn = init_db(SCHEMA_VERSION)
    llm = initialize_llm()

    # ============================================================
    # Sidebar - Authentication
    # ============================================================
    with st.sidebar:
        st.markdown("### 🔐 Account")

        if "user_id" not in st.session_state:
            auth_tab = st.radio("", ["Login", "Register"], horizontal=True, label_visibility="collapsed")
            st.divider()

            if auth_tab == "Login":
                st.markdown("#### Sign In")
                username = st.text_input("Username or Email", key="login_user", placeholder="Enter username or email")
                show_password = st.checkbox("Show password", key="show_pw_login")
                password = st.text_input("Password", type="text" if show_password else "password", key="login_pass", placeholder="Enter password")

                if st.button("🔑 Login", use_container_width=True):
                    ok, msg, uid, uname = login_user(conn, username, password)
                    if ok:
                        st.session_state["user_id"] = uid
                        st.session_state["username"] = uname
                        st.success(msg)
                        time.sleep(0.5)
                        st.rerun()
                    else:
                        st.error(msg)

            else:
                st.markdown("#### Create Account")
                new_username = st.text_input("Username", key="reg_user", placeholder="Choose a username")
                new_email = st.text_input("Email", key="reg_email", placeholder="your.email@example.com")
                show_password2 = st.checkbox("Show password", key="show_pw_reg")
                new_password = st.text_input("Password", type="text" if show_password2 else "password", key="reg_pass", placeholder="Min 6 chars, include number")
                confirm_password = st.text_input("Confirm Password", type="text" if show_password2 else "password", key="reg_pass2", placeholder="Re-enter password")

                if st.button("✨ Create Account", use_container_width=True):
                    if new_password != confirm_password:
                        st.error("Passwords do not match.")
                    else:
                        ok, msg = create_user(conn, new_username, new_email, new_password)
                        if ok:
                            st.success(msg)
                        else:
                            st.error(msg)

        else:
            st.success(f"👤 **{st.session_state.get('username', '')}**")
            if st.button("🚪 Logout", use_container_width=True):
                st.session_state.clear()
                st.rerun()

    # Show hero page if not logged in
    if "user_id" not in st.session_state:
        render_hero()
        render_features()
        st.divider()
        st.info("👈 Please login or create an account to continue")
        return

    # Check if LLM initialized
    if llm is None:
        st.error("⚠️ AI model not initialized. Please check your API key configuration.")
        st.stop()

    user_id = st.session_state["user_id"]

    # ============================================================
    # Sidebar - Module Management
    # ============================================================
    with st.sidebar:
        st.divider()
        st.markdown("### 📚 Your Modules")
        
        new_module = st.text_input("", placeholder="e.g., Biology, Python...", key="new_module_input", label_visibility="collapsed")
        if st.button("➕ Create / Open Module", use_container_width=True):
            if new_module.strip():
                module_id = get_or_create_module(conn, user_id, new_module.strip())
                st.session_state["module_id"] = module_id
                st.session_state["module_name"] = new_module.strip()
                st.success(f"📂 Opened: {new_module.strip()}")
                time.sleep(0.3)
                st.rerun()
            else:
                st.warning("Please enter a module name")

        st.divider()
        modules = list_modules(conn, user_id)
        if modules:
            st.markdown("**Your Modules:**")
            for mod_id, mod_name in modules[:15]:
                if st.button(f"📂 {mod_name}", key=f"open_mod_{mod_id}", use_container_width=True):
                    st.session_state["module_id"] = mod_id
                    st.session_state["module_name"] = mod_name
                    st.rerun()

    # Check if module is selected
    if "module_id" not in st.session_state:
        render_hero()
        st.info("📚 Create or select a module from the sidebar to get started")
        return

    module_id = st.session_state["module_id"]
    module_name = st.session_state.get("module_name", "Unknown Module")

    # Module header
    st.markdown(f'<h2 style="color:#22d3ee; font-weight:800;">📖 {module_name}</h2>', unsafe_allow_html=True)
    
    # Navigation buttons
    col1, col2, col3 = st.columns(3)
    with col1:
        if st.button("📝 Notes → Quiz", key="nav_notes", use_container_width=True):
            st.session_state["active_tab"] = "notes"
            st.rerun()
    with col2:
        if st.button("✅ Take Quiz", key="nav_take", use_container_width=True):
            st.session_state["active_tab"] = "take"
            st.rerun()
    with col3:
        if st.button("📊 Progress", key="nav_progress", use_container_width=True):
            st.session_state["active_tab"] = "progress"
            st.rerun()
    
    st.divider()

    # ============================================================
    # TAB 1: Notes → Quiz Generation
    # ============================================================
    if st.session_state["active_tab"] == "notes":
        st.markdown('<div class="qc-card">', unsafe_allow_html=True)
        st.markdown("### 📄 Upload Your Study Notes")
        st.markdown('<p style="color: var(--muted);">Supported: PDF, DOCX, RTF, ODT, PPTX, TXT</p>', unsafe_allow_html=True)

        uploaded_file = st.file_uploader(
            "Choose file",
            type=["pdf", "docx", "rtf", "odt", "pptx", "txt"],
            label_visibility="collapsed"
        )
        st.markdown('</div>', unsafe_allow_html=True)

        # Quiz settings
        settings_col1, settings_col2 = st.columns(2)
        with settings_col1:
            num_questions = st.number_input("📊 Number of Questions", min_value=1, max_value=50, value=10, step=1)
        with settings_col2:
            difficulty = st.selectbox("🎯 Difficulty", ["Easy", "Medium", "Hard"], index=1)

        if uploaded_file:
            with st.spinner("📖 Extracting text..."):
                notes_text, error = extract_notes(uploaded_file)

            if error:
                st.error(f"❌ {error}")
                st.stop()

            if not notes_text:
                st.error("❌ No text could be extracted. Please check your file.")
                st.stop()

            notes_hash = compute_hash(notes_text)

            # Show preview
            st.markdown('<div class="qc-card">', unsafe_allow_html=True)
            st.markdown("#### 📄 Notes Preview")
            preview = notes_text[:1500] + ("..." if len(notes_text) > 1500 else "")
            st.text_area("", preview, height=250, disabled=True, label_visibility="collapsed")
            st.caption(f"📏 Total: {len(notes_text):,} characters")
            st.markdown('</div>', unsafe_allow_html=True)

            # Action buttons
            action_col1, action_col2 = st.columns([3, 1])
            with action_col1:
                generate_btn = st.button("✨ Generate Quiz with AI", use_container_width=True, type="primary")
            with action_col2:
                clear_btn = st.button("🧹 Clear", use_container_width=True)

            if clear_btn:
                if "selected_quiz_id" in st.session_state:
                    del st.session_state["selected_quiz_id"]
                st.success("Cleared")
                st.rerun()

            if generate_btn:
                progress = st.progress(0)
                status = st.empty()
                
                try:
                    status.text("🤖 AI analyzing your notes...")
                    progress.progress(25)
                    
                    quiz_data = generate_quiz_with_retry(llm, notes_text, num_questions, difficulty, max_attempts=3)
                    progress.progress(75)
                    
                    status.text("💾 Saving to database...")
                    quiz_id = save_quiz(conn, user_id, module_id, notes_hash, quiz_data, difficulty)
                    progress.progress(100)
                    
                    st.session_state["selected_quiz_id"] = quiz_id
                    status.empty()
                    progress.empty()
                    
                    st.success(f"✅ Quiz #{quiz_id} created successfully!")
                    
                    # Next action buttons
                    next_col1, next_col2, next_col3 = st.columns(3)
                    with next_col1:
                        if st.button("➡️ Take Quiz Now", use_container_width=True):
                            st.session_state["active_tab"] = "take"
                            st.rerun()
                    with next_col2:
                        if st.button("📊 View Progress", use_container_width=True):
                            st.session_state["active_tab"] = "progress"
                            st.rerun()
                    with next_col3:
                        st.download_button(
                            "⬇️ Download JSON",
                            data=json.dumps(quiz_data, indent=2),
                            file_name=f"quiz_{quiz_id}.json",
                            mime="application/json",
                            use_container_width=True
                        )
                
                except Exception as e:
                    progress.empty()
                    status.empty()
                    st.error(f"❌ Quiz generation failed: {e}")
                    st.info("💡 Try: shorter notes, fewer questions, or different difficulty")
        
        else:
            st.info("👆 Upload your study notes to generate a quiz")

    # ============================================================
    # TAB 2: Take Quiz
    # ============================================================
    elif st.session_state["active_tab"] == "take":
        quiz_list = list_quiz_ids_for_module(conn, user_id, module_id, limit=50)

        if not quiz_list:
            st.info("📝 No quizzes available. Generate one from 'Notes → Quiz' tab.")
        else:
            # Quiz selector
            options = [(qid, f"Quiz #{qid} • {ts} • {diff}") for (qid, ts, diff) in quiz_list]
            default_qid = st.session_state.get("selected_quiz_id", options[0][0])

            selected_label = st.selectbox(
                "📋 Select Quiz",
                options=[label for _, label in options],
                index=next((i for i, (qid, _) in enumerate(options) if qid == default_qid), 0)
            )

            selected_quiz_id = next(qid for qid, label in options if label == selected_label)
            st.session_state["selected_quiz_id"] = selected_quiz_id

            # Action buttons
            top_col1, top_col2 = st.columns(2)
            with top_col1:
                if st.button("🔄 Retake (Clear Answers)", use_container_width=True):
                    quiz_to_clear = load_quiz(conn, selected_quiz_id)
                    if quiz_to_clear:
                        for q in quiz_to_clear["questions"]:
                            key = f"q_{selected_quiz_id}_{q['question_id']}"
                            if key in st.session_state:
                                del st.session_state[key]
                    st.success("✅ Answers cleared")
                    st.rerun()
            with top_col2:
                if st.button("📝 Back to Notes", use_container_width=True):
                    st.session_state["active_tab"] = "notes"
                    st.rerun()

            # Load and display quiz
            quiz = load_quiz(conn, selected_quiz_id)
            if not quiz:
                st.error("Quiz not found")
                st.stop()

            st.markdown(f"**Quiz ID:** {quiz['quiz_id']} • **Difficulty:** {quiz['difficulty']} • **Created:** {quiz['timestamp']}")
            st.divider()

            # Quiz form
            user_choices = {}
            with st.form(key=f"quiz_form_{selected_quiz_id}"):
                for idx, question in enumerate(quiz["questions"], 1):
                    st.markdown('<div class="quiz-question">', unsafe_allow_html=True)
                    st.markdown(f"### Question {idx}")
                    st.markdown(f"**{question['question']}**")
                    
                    choice = st.radio(
                        "Your answer:",
                        options=[0, 1, 2, 3],
                        format_func=lambda x: question["options"][x],
                        key=f"q_{selected_quiz_id}_{question['question_id']}",
                        index=None,
                        label_visibility="collapsed"
                    )
                    user_choices[question["question_id"]] = choice
                    st.markdown('</div>', unsafe_allow_html=True)

                st.divider()
                submitted = st.form_submit_button("✅ Submit Quiz", use_container_width=True, type="primary")

            # Process submission
            if submitted:
                missing = [qid for qid, ans in user_choices.items() if ans is None]
                
                if missing:
                    st.error(f"⚠️ Please answer all questions. Missing: {len(missing)}")
                else:
                    score, total, percentage = save_user_answers_and_performance(conn, user_id, module_id, quiz, user_choices)

                    # Show result
                    if percentage >= 90:
                        emoji, message = "🏆", "Outstanding!"
                    elif percentage >= 75:
                        emoji, message = "🎯", "Great job!"
                    elif percentage >= 60:
                        emoji, message = "👍", "Good effort!"
                    else:
                        emoji, message = "💪", "Keep practicing!"

                    st.success(f"{emoji} {message} Your score: **{score}/{total}** ({percentage:.1f}%)")
                    st.divider()

                    # Review answers
                    st.markdown("### 📋 Review Your Answers")
                    for idx, question in enumerate(quiz["questions"], 1):
                        chosen = int(user_choices[question["question_id"]])
                        correct = int(question["answer"])

                        if chosen == correct:
                            st.success(f"**Q{idx}:** ✅ Correct - {question['options'][chosen]}")
                        else:
                            st.error(f"**Q{idx}:** ❌ Incorrect")
                            st.markdown(f"Your answer: **{question['options'][chosen]}**")
                            st.markdown(f"Correct answer: **{question['options'][correct]}**")

                        if question.get("explanation"):
                            with st.expander("💡 Explanation"):
                                st.write(question["explanation"])

                    # Next actions
                    result_col1, result_col2 = st.columns(2)
                    with result_col1:
                        if st.button("📊 View Progress", key="view_prog", use_container_width=True):
                            st.session_state["active_tab"] = "progress"
                            st.rerun()
                    with result_col2:
                        if st.button("📝 New Quiz", key="new_quiz", use_container_width=True):
                            st.session_state["active_tab"] = "notes"
                            st.rerun()

    # ============================================================
    # TAB 3: Progress Analytics
    # ============================================================
    else:
        performance_df = fetch_performance(conn, user_id, module_id)

        if performance_df.empty:
            st.info("📊 No performance data yet. Complete a quiz to see your progress!")
        else:
            # Stats cards
            stat_col1, stat_col2, stat_col3 = st.columns(3)
            with stat_col1:
                st.markdown('<div class="qc-card">', unsafe_allow_html=True)
                st.metric("Total Attempts", len(performance_df))
                st.markdown('</div>', unsafe_allow_html=True)
            with stat_col2:
                st.markdown('<div class="qc-card">', unsafe_allow_html=True)
                st.metric("Average Score", f"{performance_df['percentage'].mean():.1f}%")
                st.markdown('</div>', unsafe_allow_html=True)
            with stat_col3:
                st.markdown('<div class="qc-card">', unsafe_allow_html=True)
                st.metric("Best Score", f"{performance_df['percentage'].max():.1f}%")
                st.markdown('</div>', unsafe_allow_html=True)

            st.divider()

            # Performance chart
            st.markdown("### 📈 Performance Trend")
            fig = plot_progress(performance_df)
            st.pyplot(fig, use_container_width=True)

            st.divider()

            # Detailed history
            st.markdown("### 📋 Quiz History")
            display_df = performance_df.copy()
            display_df["timestamp"] = display_df["timestamp"].dt.strftime("%Y-%m-%d %H:%M")
            display_df = display_df.rename(columns={
                "timestamp": "Date & Time",
                "score": "Score",
                "total": "Total",
                "percentage": "Percentage (%)"
            })
            st.dataframe(display_df, use_container_width=True, hide_index=True)


if __name__ == "__main__":
    main()